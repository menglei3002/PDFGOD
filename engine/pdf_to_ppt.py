"""PDF to PPT conversion: render PDF pages as images, embed in PPTX."""

import os
import tempfile
from pathlib import Path


def convert(input_path: str, output_path: str, report_fn) -> str:
    report_fn(10, "Rendering PDF pages to images...")

    # Render PDF to images first
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    doc = fitz.open(input_path)
    total = len(doc)
    prs = Presentation()

    # Default slide size
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    with tempfile.TemporaryDirectory() as tmpdir:
        for i in range(total):
            page = doc[i]
            # Render at decent resolution
            mat = fitz.Matrix(2.0, 2.0)  # ~144 DPI
            pix = page.get_pixmap(matrix=mat)
            img_path = os.path.join(tmpdir, f"page_{i:03d}.png")
            pix.save(img_path)

            # Add slide with the image
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0),
                                     slide_width, slide_height)

            pct = 10 + int(80 * (i + 1) / total)
            report_fn(pct, f"Building slide {i+1}/{total}...")

    doc.close()

    report_fn(90, "Saving PPTX file...")
    prs.save(output_path)

    if os.path.exists(output_path):
        return os.path.abspath(output_path)
    raise RuntimeError("PPT conversion failed")
